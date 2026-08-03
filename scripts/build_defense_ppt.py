#!/usr/bin/env python3
"""
Build defense PPT for 荧光压裂液硕士答辩.
Template style from 郝乐乐开题答辩ppt.pptx, content from 论文初稿_v2.docx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import os

# ── Constants ──
DARK_BLUE = RGBColor(0x00, 0x20, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MEDIUM_BLUE = RGBColor(0x00, 0x40, 0xA0)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x1F)

SLIDE_W = Emu(9144000)  # 10 inches (16:9)
SLIDE_H = Emu(5143500)  # 5.625 inches

# ── Helper functions ──
def add_blank_slide(prs):
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(18),
                bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    # Set East Asian font
    for run in p.runs:
        run.font.name = font_name
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=Pt(18),
                          bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                          line_spacing=Pt(30), font_name='Microsoft YaHei',
                          anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple paragraphs (one per line)."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, str):
            text = line_data
            fs = font_size
            b = bold
            c = color
        else:
            text = line_data.get('text', '')
            fs = line_data.get('font_size', font_size)
            b = line_data.get('bold', bold)
            c = line_data.get('color', color)

        p.text = text
        p.font.size = fs
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.name = font_name

    return txBox

def add_section_title(slide, text, left=396575, top=85235, width=8141193, height=481863,
                      font_size=Pt(20), color=DARK_BLUE):
    """Add a section title at the top of a content slide."""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.name = 'Microsoft YaHei'
    return txBox

def add_slide_number(slide, number):
    """Add slide number in bottom right."""
    add_textbox(slide, 8500000, 4900000, 500000, 200000,
                str(number), Pt(10), False, MEDIUM_BLUE, PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    """Add a thin decorative bar at bottom."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(5070000), Emu(9144000), Emu(73500)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

def add_top_accent_line(slide):
    """Add thin accent line under section title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(396575), Emu(550000), Emu(8141193), Emu(20000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MEDIUM_BLUE
    shape.line.fill.background()

def add_toc_slide(prs, items, highlight_idx=None):
    """Add a table-of-contents / section divider slide."""
    slide = add_blank_slide(prs)
    add_textbox(slide, 1997426, 70000, 5179219, 500000,
                '目  录', Pt(36), True, DARK_BLUE, PP_ALIGN.CENTER)
    add_bottom_bar(slide)

    lines = []
    for i, item in enumerate(items):
        if highlight_idx is not None and i == highlight_idx:
            lines.append({'text': item, 'font_size': Pt(24), 'bold': True, 'color': ACCENT_GOLD})
        else:
            lines.append({'text': item, 'font_size': Pt(24), 'bold': True, 'color': DARK_BLUE})

    add_multiline_textbox(slide, 1800000, 1000000, 6000000, 3600000,
                          lines, Pt(24), True, DARK_BLUE, PP_ALIGN.LEFT,
                          line_spacing=Pt(42))
    return slide

def add_decorated_header(slide, text, subtitle=None):
    """Add a slide with centered header and decorative elements."""
    add_bottom_bar(slide)
    # Top accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), Emu(9144000), Emu(60000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    add_textbox(slide, 500000, 200000, 8144000, 700000,
                text, Pt(28), True, DARK_BLUE, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 500000, 850000, 8144000, 400000,
                    subtitle, Pt(16), False, MEDIUM_BLUE, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# BUILD DECK
# ════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── TOC items ──
toc_items = [
    '一、选题依据及意义',
    '二、国内外研究现状及存在问题',
    '三、研究目标与技术路线',
    '四、荧光粉基础物性与表面改性',
    '五、荧光压裂液体系构建与性能评价',
    '六、裂缝壁面吸附与动态驱替验证',
    '七、现场施工工艺与经济性评估',
    '八、创新点与结论',
]

# ═══════════════════════════════
# SLIDE 1: 标题页
# ═══════════════════════════════
slide = add_blank_slide(prs)
add_bottom_bar(slide)
# Top accent
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(9144000), Emu(80000))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()

add_textbox(slide, 183272, 800000, 8777456, 1200000,
            '用于压裂裂缝监测的\n荧光压裂液体系构建与性能研究',
            Pt(32), True, DARK_BLUE, PP_ALIGN.CENTER)

add_textbox(slide, 2200000, 2400000, 4800000, 900000,
            '培养单位：成都理工大学能源学院\n专    业：石油与天然气工程\n导    师：李娜\n研 究 生：郝乐乐',
            Pt(18), False, DARK_BLUE, PP_ALIGN.LEFT)

add_textbox(slide, 2500000, 3500000, 4000000, 400000,
            '硕士学位论文答辩', Pt(20), True, MEDIUM_BLUE, PP_ALIGN.CENTER)

# ═══════════════════════════════
# SLIDE 2: 答辩信息
# ═══════════════════════════════
slide = add_blank_slide(prs)
add_bottom_bar(slide)
add_textbox(slide, 2895600, 46791, 3352800, 500000,
            '答辩信息', Pt(36), True, DARK_BLUE, PP_ALIGN.CENTER)
add_top_accent_line(slide)

info_lines = [
    {'text': '论文题目：用于压裂裂缝监测的荧光压裂液体系构建与性能研究', 'font_size': Pt(22), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(14), 'bold': False, 'color': DARK_BLUE},
    {'text': '培养单位：成都理工大学能源学院', 'font_size': Pt(22), 'bold': False, 'color': DARK_BLUE},
    {'text': '专    业：石油与天然气工程（油气田开发方向）', 'font_size': Pt(22), 'bold': False, 'color': DARK_BLUE},
    {'text': '导    师：李娜', 'font_size': Pt(22), 'bold': False, 'color': DARK_BLUE},
    {'text': '研 究 生：郝乐乐', 'font_size': Pt(22), 'bold': False, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(14), 'bold': False, 'color': DARK_BLUE},
    {'text': '已修学分：已修够学分  |  学位英语：80分', 'font_size': Pt(20), 'bold': False, 'color': MEDIUM_BLUE},
    {'text': '研究方向：油田化学', 'font_size': Pt(20), 'bold': False, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 1500000, 900000, 6200000, 3700000,
                      info_lines, Pt(22), False, DARK_BLUE, PP_ALIGN.LEFT, Pt(36))

# ═══════════════════════════════
# SLIDE 3: 目录
# ═══════════════════════════════
add_toc_slide(prs, toc_items)

# ═══════════════════════════════
# SLIDE 4-5: 选题依据及意义
# ═══════════════════════════════
# --- Slide 4: 研究背景 ---
slide = add_blank_slide(prs)
add_section_title(slide, '一、选题依据及意义——研究背景')
add_top_accent_line(slide)
add_bottom_bar(slide)

bg_lines = [
    {'text': '▎ 水力压裂是非常规油气资源商业化开发的核心技术', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 全球已完成逾250万次压裂施工，累计为美国新增可采石油储量约30%、天然气约90%', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '• 裂缝延伸范围直接决定储层改造体积（SRV）——页岩油气产能的关键控制参数', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(10), 'bold': False, 'color': BLACK},
    {'text': '▎ 现有裂缝监测技术的结构性缺陷：缺乏原位实物证据', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 微地震、DAS/DTS、示踪剂等均属间接测量，依赖反演，多解性强', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '• 压后取心作为最直接验证手段，现有技术难以在岩心上留下可辨识标记', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '• 裂缝诊断结果始终缺乏"可取证"的实物依据', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(10), 'bold': False, 'color': BLACK},
    {'text': '▎ 荧光示踪技术：弥补"信号-实物"鸿沟的可行路径', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 向压裂液中引入荧光标记材料，压后取心在紫外下直接观察荧光分布', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '• 已有探索：Ishida等(2025)荧光树脂、Guryanov等(2019)量子点微球', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '• 核心工程问题：如何在常规胍胶压裂液内实现荧光示踪材料的稳定输送与壁面锚定？', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 650000, 8200000, 4200000,
                      bg_lines, Pt(16), False, BLACK, PP_ALIGN.LEFT, Pt(26))

# --- Slide 5: 技术方案对比 ---
slide = add_blank_slide(prs)
add_section_title(slide, '一、选题依据及意义——技术方案对比与定位')
add_top_accent_line(slide)
add_bottom_bar(slide)

compare_text = (
    '本研究以 SrAl₂O₄:Eu²⁺,Dy³⁺ 无机长余辉荧光粉为示踪介质、\n'
    'HPG胍胶冻胶为压裂液载体，构建与现有压裂工艺完全兼容的\n'
    '压裂液波及范围可视化方法。\n\n'
    '该无机荧光粉的核心优势在于：\n'
    '发光中心（Eu²⁺）受刚性晶格保护，对过硫酸铵破胶剂的\n'
    '氧化环境具有本征化学惰性——确保示踪材料在破胶返排\n'
    '全过程中保持荧光信号完整性。\n\n'
    '技术定位：微地震/DAS等间接监测方法的辅助校准工具\n'
    '——通过压后取心实物证据验证间接反演结果的可靠性。'
)
add_textbox(slide, 500000, 600000, 4200000, 4200000,
            compare_text, Pt(16), False, BLACK, PP_ALIGN.LEFT)

# Table: comparison of monitoring technologies
table_data = [
    ['方法', '测量方式', '原位证据', '微裂缝', '定量能力'],
    ['微地震', '间接反演', '否', '否', '否'],
    ['DAS/DTS', '间接反演', '否', '部分', '半定量'],
    ['化学示踪剂', '返排分析', '否', '否', '半定量'],
    ['放射性示踪', '返排分析', '否', '否', '否'],
    ['荧光树脂', '直接成像', '是', '差', '否'],
    ['本研究', '取心+紫外成像', '是', '可', '半定量'],
]
rows, cols = len(table_data), len(table_data[0])
table_shape = slide.shapes.add_table(rows, cols,
    Emu(5000000), Emu(500000), Emu(3800000), Emu(4200000))
table = table_shape.table

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = 'Microsoft YaHei'
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r == rows - 1:
                p.font.bold = True
                p.font.color.rgb = MEDIUM_BLUE
            else:
                p.font.color.rgb = BLACK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif r == rows - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)

# ═══════════════════════════════
# SLIDE 6: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 1)

# ═══════════════════════════════
# SLIDE 7-8: 国内外研究现状与存在问题
# ═══════════════════════════════
# --- Slide 7: 研究现状 ---
slide = add_blank_slide(prs)
add_section_title(slide, '二、国内外研究现状')
add_top_accent_line(slide)
add_bottom_bar(slide)

status_lines = [
    {'text': '▎ 裂缝监测技术', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• DAS/DTS：Molenaar(2012)首次井下应用 → Jin & Roy(2017)低频DAS-裂缝几何定量关系', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 微地震：Maxwell(2002)Barnett页岩裂缝网络成像 → SRV概念(Mayerhofer 2010)', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 示踪剂：化学示踪剂返排分析(Salman 2014) → 量子点微球GeoSplit(Guryanov 2019)', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 共性局限：间接反演多解性强，裂缝真实空间展布难以唯一确定', 'font_size': Pt(14), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 荧光示踪裂缝可视化', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 直接可视化路线：荧光树脂(Chen 2014, Takeuchi 2025)——与常规压裂工艺不兼容', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 返排分析路线：量子点、碳点、介孔硅纳米颗粒——依赖间接信号，无法在壁面留下实物标记', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 稀土铝酸盐荧光材料', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• SrAl₂O₄:Eu²⁺,Dy³⁺：Matsuzawa(1996)首次报道 → 余辉亮度比ZnS:Cu,Co高10倍以上', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 表面包覆：SiO₂(Qi 2017)、ALD Al₂O₃/TiO₂(Karacaoglu 2020)、KH570桥联(Lyu 2020)', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 研究空白：所有包覆方案均面向涂料/塑料行业，压裂液工程环境的系统验证尚属空白', 'font_size': Pt(14), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 600000, 8200000, 4200000,
                      status_lines, Pt(14), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 8: 存在问题 ---
slide = add_blank_slide(prs)
add_section_title(slide, '二、存在的主要问题与研究空白')
add_top_accent_line(slide)
add_bottom_bar(slide)

problem_lines = [
    {'text': '当前存在三个层次的技术缺口：', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '① 材料层面——分散悬浮稳定性', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  高密度荧光粉(3.6~4.0 g/cm³)在HPG基液(~1.0 g/cm³)中悬浮稳定性未解决；', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '  荧光粉溶出多价阳离子(Al³⁺/Ca²⁺)可能干扰胍胶交联反应', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '② 工艺层面——"分散-吸附"功能时序切换', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  注入期需稳定悬浮（亲水分散），关井后需壁面锚定（化学结合），需求截然相反；', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '  缺乏利用压裂施工自身时序（注入→关井→返排）驱动功能切换的工程化方案', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '③ 评价层面——动态全流程验证缺失', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  已有研究多采用静态浸泡/简单注入后成像，跳过了注入剪切、关井破胶和返排冲刷三个工程必经环节；', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '  动态驱替模拟"注入—破胶—返排"全过程条件下的荧光示踪效果尚未被验证', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 400000, 550000, 8400000, 4300000,
                      problem_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(26))

# ═══════════════════════════════
# SLIDE 9: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 2)

# ═══════════════════════════════
# SLIDE 10: 研究目标与技术路线
# ═══════════════════════════════
slide = add_blank_slide(prs)
add_section_title(slide, '三、研究目标与技术路线')
add_top_accent_line(slide)
add_bottom_bar(slide)

goal_lines = [
    {'text': '▎ 研究目标', 'font_size': Pt(22), 'bold': True, 'color': DARK_BLUE},
    {'text': '以SrAl₂O₄:Eu²⁺,Dy³⁺为示踪介质、HPG压裂液为载体，以压后取心紫外直接观察为检测手段，', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '构建一种与现有压裂工艺完全兼容的压裂液波及范围实物验证方法。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 研究内容（六部分）', 'font_size': Pt(22), 'bold': True, 'color': DARK_BLUE},
    {'text': '（1）荧光粉基础物性表征与储层环境适应性评价', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '（2）KH550+PEG4000双层表面改性工艺优化与协同分散机理', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '（3）高浓度荧光母液研制与压裂液体系标准化性能评价', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '（4）"分散-吸附"功能切换机制与砂岩壁面吸附规律', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '（5）动态驱替模拟压裂-返排全过程与荧光信号半定量分析', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '（6）"母液预配+在线稀释"现场施工工艺方案与经济性评估', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 主线：材料改性 → 体系构建 → 动态验证 → 工程转化', 'font_size': Pt(18), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 600000, 8200000, 4000000,
                      goal_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(24))

# --- Slide 11: 技术路线图（文字版）---
slide = add_blank_slide(prs)
add_section_title(slide, '三、技术路线图')
add_top_accent_line(slide)
add_bottom_bar(slide)

# Create a visual flow chart with text boxes and arrows
route_items = [
    ('染料类型优选\n（稀土铝酸盐荧光粉）', 300000, 700000),
    ('荧光粉物性研究\n（热/氧化/耐盐/pH）', 300000, 1800000),
    ('颗粒表面修饰\n（KH550+PEG4000双层）', 300000, 2900000),
    ('荧光压裂液研究', 5200000, 500000),
    ('  └─ 母液配制（≥40 g/L）', 5200000, 1300000),
    ('  └─ 基液性能评价', 5200000, 2000000),
    ('  └─ 交联冻胶/悬砂/破胶', 5200000, 2700000),
    ('岩心吸附研究', 5200000, 3600000),
    ('  ├─ 静态吸附（热力学/动力学）', 5200000, 3900000),
    ('  └─ 动态驱替（五步串联流程）', 5200000, 4200000),
]

for text, left, top in route_items:
    fs = Pt(13) if text.startswith('  ') else Pt(14)
    b = not text.startswith('  ')
    add_textbox(slide, left, top, 3500000, 400000,
                text.strip(), fs, b, DARK_BLUE if b else BLACK, PP_ALIGN.LEFT)

# Right side: outcome boxes
outcomes = [
    ('判读方法建立', 300000, 4000000),
    ('施工方案设计\n（母液预配+在线稀释）', 5200000, 4500000),
]
for text, left, top in outcomes:
    add_textbox(slide, left, top, 3500000, 500000,
                text, Pt(14), True, MEDIUM_BLUE, PP_ALIGN.LEFT)

# Arrows (using simple text shapes)
add_textbox(slide, 3850000, 750000, 1300000, 400000,
            '───────▶', Pt(24), True, ACCENT_GOLD, PP_ALIGN.CENTER)
add_textbox(slide, 3850000, 1850000, 1300000, 400000,
            '───────▶', Pt(24), True, ACCENT_GOLD, PP_ALIGN.CENTER)
add_textbox(slide, 3850000, 2950000, 1300000, 400000,
            '───────▶', Pt(24), True, ACCENT_GOLD, PP_ALIGN.CENTER)

# Bottom: evaluation standards
add_textbox(slide, 1000000, 4650000, 7000000, 400000,
            '性能评价依据：SY/T 6376-2008《压裂液通用技术条件》| SY/T 5107-2016《水基压裂液性能评价方法》',
            Pt(11), False, MEDIUM_BLUE, PP_ALIGN.CENTER)

# ═══════════════════════════════
# SLIDE 12: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 3)

# ═══════════════════════════════
# SLIDE 13-15: 荧光粉基础物性与表面改性
# ═══════════════════════════════
# --- Slide 13a: 荧光粉表征 ---
slide = add_blank_slide(prs)
add_section_title(slide, '四、荧光粉基础物性——材料表征与环境适应性')
add_top_accent_line(slide)
add_bottom_bar(slide)

char_lines = [
    {'text': '▎ 基础物性', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 材料：商用SrAl₂O₄:Eu²⁺,Dy³⁺（固相法合成，1000目，D50≈13μm）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 表征：XRD晶相分析 | SEM-EDS形貌与元素 | 激光粒度分布 | 荧光激发/发射光谱+余辉衰减曲线', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 储层环境适应性评价（实验数据待补充）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 热稳定性：60~150°C，0~168h → 获取发光衰减速率常数k和半衰期t₁/₂', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 化学稳定性：NaCl 0~100 g/L，90°C → 同步监测pH以评估水解程度', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 氧化稳定性：过硫酸铵0~0.2% w/v，90°C → 关键指标——与有机荧光材料的本质区别', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• pH适应性：pH 3~11缓冲溶液，90°C', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(8), 'bold': False, 'color': BLACK},
    {'text': '▎ 拟合模型：一级衰减动力学 ln(I/I₀) = -kt（每项实验）', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 600000, 8200000, 4200000,
                      char_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(24))

# --- Slide 14: 双层表面改性 ---
slide = add_blank_slide(prs)
add_section_title(slide, '四、双层表面改性方案设计')
add_top_accent_line(slide)
add_bottom_bar(slide)

mod_lines = [
    {'text': '▎ 改性目标：构建兼具三重功能的双层结构', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '  内层 — KH550硅烷偶联剂', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '    • 化学锚固：Si-O-Al共价键锚定荧光粉表面', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '    • 功能基团：末端-NH₂（氨基）——为岩石壁面Si-OH提供吸附活性位点', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '    • 耐水解屏障：致密硅氧烷网络阻隔水分子渗透', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '  外层 — PEG4000（聚乙二醇，Mw=4000）', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '    • 物理屏蔽：长链构象熵+渗透排斥 → 空间位阻分散', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '    • 可脱附设计：高温+氧化条件下PEG链从颗粒表面脱附', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '    • 暴露活性位点：脱附后释放KH550氨基 → 与岩石Si-OH锚定', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 协同分散体系：柠檬酸（螯合屏蔽Al³⁺）+ Triton X-100（界面润湿）+ PEG（空间位阻）', 'font_size': Pt(16), 'bold': True, 'color': DARK_BLUE},
    {'text': '▎ 工艺优化：L9(3⁴)正交实验——KH550用量×PEG分子量×PEG浓度×搅拌时间', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4300000,
                      mod_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 15: 分散稳定性 ---
slide = add_blank_slide(prs)
add_section_title(slide, '四、改性粉在胍胶基液中的分散稳定性')
add_top_accent_line(slide)
add_bottom_bar(slide)

disp_lines = [
    {'text': '▎ 六组对照实验设计——量化各组分独立贡献与协同效应', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '组1：双层改性粉 + HPG基液（空白对照）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '组2：空白 + 0.1 wt%柠檬酸（螯合Al³⁺）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '组3：空白 + 0.05 wt% Triton X-100（界面润湿）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '组4：空白 + 5 wt% PEG4000（额外空间位阻）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '组5：空白 + 柠檬酸 + Triton（两两组合）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '组6：空白 + 柠檬酸 + Triton + PEG（三者全加）', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 评估指标', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 静态沉降曲线：相对浊度保持率 vs 时间（0~120 min）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 激光粒度：D50和Span值随时间变化（反映团聚动力学）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• Zeta电位-pH曲线（颗粒表面电荷稳定性）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 协同效应量化：协同项 = 组6效果 - (组2+组3+组4 - 3×组1)', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '▎ 交联影响验证：含/不含荧光粉的HPG冻胶交联时间、G\'/G"粘弹性、SEM微观结构对比', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4300000,
                      disp_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# ═══════════════════════════════
# SLIDE 16: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 4)

# ═══════════════════════════════
# SLIDE 17-18: 荧光压裂液体系构建与性能评价
# ═══════════════════════════════
# --- Slide 17: 母液+基液 ---
slide = add_blank_slide(prs)
add_section_title(slide, '五、荧光压裂液体系构建——母液制备与基液评价')
add_top_accent_line(slide)
add_bottom_bar(slide)

sys_lines = [
    {'text': '▎ 高浓度荧光母液制备', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 浓度目标：40~60 g/L（改性荧光粉 + 柠檬酸 + Triton X-100 + 去离子水）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 工艺：高速均质机10000~15000 rpm × 5 min → 超声40 kHz × 10 min脱泡', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 质量指标：7天静置分层≤5% | 离心浊度比TR≥0.90 | 荧光强度保持率≥90%', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 荧光压裂液终液', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 配方：0.5 wt% HPG基液 + 0.5% v/v荧光母液 → 复配得终液', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 基液评价（SY/T 5107-2016）：表观粘度（170 s⁻¹, 25°C）、pH、密度', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 与空白HPG基液对比 → 验证荧光母液不显著改变基液基础物理性质', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 交联冻胶性能（SY/T 5107-2016）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 交联剂：有机硼延缓交联剂0.3% v/v', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 耐温耐剪切：25→120°C，170 s⁻¹恒剪 → 表观粘度≥50 mPa·s维持60 min以上', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 粘弹性：振荡频率扫描0.1~10 Hz → G\' > G"弹性主导准则（支撑剂悬浮能力）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4300000,
                      sys_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 18: 悬砂+破胶+伤害 ---
slide = add_blank_slide(prs)
add_section_title(slide, '五、压裂液性能评价——悬砂/破胶/地层伤害')
add_top_accent_line(slide)
add_bottom_bar(slide)

perf_lines = [
    {'text': '▎ 悬砂性能（SY/T 6376-2008）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 20/40目陶粒（~2.7 g/cm³, 480 kg/m³），室温/储层温度', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 标准：储层温度下沉降速度 ≤ 0.5 mm/min', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 破胶性能（SY/T 5107-2016）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 破胶剂：过硫酸铵0.05~0.2% w/v，60~120°C → 破胶至粘度≤10 mPa·s', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 残渣含量：过滤-干燥-称重法（mg/L），与空白HPG冻胶对比', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 地层伤害评价', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 滤饼伤害：HPHT滤失仪，3.5 MPa / 90°C → 滤饼厚度+SEM形貌+EDS分布', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 岩心驱替伤害（SY/T 6540-2002）：伤害前后煤油渗透率对比，剥离荧光粉增量贡献ΔD', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 导流能力（SY/T 6302-2009）：API导流室，20/40目陶粒，闭合压力10~60 MPa', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '核心结论：荧光母液的引入不显著损害压裂液工程性能', 'font_size': Pt(18), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4300000,
                      perf_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# ═══════════════════════════════
# SLIDE 19: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 5)

# ═══════════════════════════════
# SLIDE 20-22: 裂缝壁面吸附与动态驱替
# ═══════════════════════════════
# --- Slide 20: 功能切换假说 ---
slide = add_blank_slide(prs)
add_section_title(slide, '六、"分散-吸附"功能切换机制')
add_top_accent_line(slide)
add_bottom_bar(slide)

switch_lines = [
    {'text': '▎ 核心工作假说：利用压裂工程时序驱动表面化学功能切换', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '注入阶段（分散态）', 'font_size': Pt(22), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  PEG4000外层充分伸展 → 空间位阻排斥 → 颗粒均匀悬浮', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '  KH550氨基被PEG屏蔽 → 不与HPG/交联剂发生副反应', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '关井破胶阶段（切换窗口）', 'font_size': Pt(22), 'bold': True, 'color': ACCENT_GOLD},
    {'text': '  过硫酸铵氧化 + 90°C高温 → PEG链从颗粒表面脱附', 'font_size': Pt(16), 'bold': True, 'color': BLACK},
    {'text': '  KH550末端-NH₂暴露 → 质子化为-NH₃⁺', 'font_size': Pt(16), 'bold': True, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '返排+取心阶段（锚定态）', 'font_size': Pt(22), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  三级协同锚定机制：', 'font_size': Pt(16), 'bold': True, 'color': BLACK},
    {'text': '  ① 静电吸附：-NH₃⁺ ⟷ 砂岩表面Si-O⁻', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '  ② 氢键：-NH₂···HO-Si≡', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '  ③ 化学缩合：-NH₂ + HO-Si → -NH-Si + H₂O（高温催化）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 化学证据（FTIR + XPS + Zeta电位交叉验证）——实验数据待补充', 'font_size': Pt(15), 'bold': False, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 550000, 8200000, 4350000,
                      switch_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(21))

# --- Slide 21: 静态吸附 ---
slide = add_blank_slide(prs)
add_section_title(slide, '六、静态吸附实验——热力学与动力学')
add_top_accent_line(slide)
add_bottom_bar(slide)

adsorp_lines = [
    {'text': '▎ 批吸附实验设计', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 砂岩薄片(10×10×2 mm) + 破胶处理荧光粉悬浮液(0.1~5.0 g/L)', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 模拟地层水(50 g/L NaCl + 2 g/L CaCl₂)，25/50/80°C恒温振荡至平衡', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• UV-Vis/荧光光谱仪测定残余浓度 → 计算单位面积吸附量(μg/cm²)', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 等温吸附模型拟合', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• Langmuir: q_e = q_max·K_L·C_e/(1+K_L·C_e)', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• Freundlich: q_e = K_F·C_e^(1/n)', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• Temkin —— 以R²和AIC准则选择最优模型', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 热力学参数（Van\'t Hoff方程）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• ΔG° = -RT·ln(K_L)：判断吸附自发性 → 预期ΔG° < 0', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• ln(K_L) vs 1/T → ΔH°（吸附热）和ΔS°（熵变）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 动力学模型', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 拟一级：ln(q_e-q_t) vs t | 拟二级：t/q_t vs t | 颗粒内扩散：q_t vs t^(1/2)', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 识别速率控制步骤（膜扩散/颗粒内扩散/表面吸附）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4350000,
                      adsorp_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 22: 动态驱替 ---
slide = add_blank_slide(prs)
add_section_title(slide, '六、动态驱替四阶段实验——工程可行性完整证据链')
add_top_accent_line(slide)
add_bottom_bar(slide)

dynamic_lines = [
    {'text': '实验装置：可调节裂缝宽度的岩心夹持系统（砂岩/巴西劈裂法造缝/0.1~2.0 mm）', 'font_size': Pt(14), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(4), 'bold': False, 'color': BLACK},
    {'text': '阶段一 ▸ 注入', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  荧光HPG压裂液+有机硼交联剂在线混配，90°C恒流注入3倍裂缝体积', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '  证据A：注入端压力-注入PV曲线 → 判断颗粒堵塞/架桥（p_in是否单调上升无平台）', 'font_size': Pt(14), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(4), 'bold': False, 'color': BLACK},
    {'text': '阶段二 ▸ 关井破胶', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  回压5 MPa，90°C密闭12 h，过硫酸铵破胶 → PEG脱附+KH550氨基暴露', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '  此为"分散→吸附"功能切换的静态窗口期，不采集数据', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(4), 'bold': False, 'color': BLACK},
    {'text': '阶段三 ▸ 返排', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  模拟地层水反向低速驱替5~10倍裂缝体积', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '  证据B：返排液荧光浓度-CPV曲线 → 计算净残留率（目标>90%）→ 验证锚定牢固度', 'font_size': Pt(14), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '', 'font_size': Pt(4), 'bold': False, 'color': BLACK},
    {'text': '阶段四 ▸ 取心+紫外成像', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '  拆卸岩心，365 nm紫外+520±10 nm带通滤光片成像（所有参数锁定不变）', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '  证据C：裂缝壁面紫外图像+ImageJ灰度分析 → 验证空间对应性', 'font_size': Pt(14), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 450000, 500000, 8300000, 4400000,
                      dynamic_lines, Pt(14), False, BLACK, PP_ALIGN.LEFT, Pt(20))

# ═══════════════════════════════
# SLIDE 23: 灰度-裂缝宽度关系
# ═══════════════════════════════
slide = add_blank_slide(prs)
add_section_title(slide, '六、荧光信号与裂缝几何的半定量关系')
add_top_accent_line(slide)
add_bottom_bar(slide)

quant_lines = [
    {'text': '▎ 分析框架：递进式——无论实验结果落在哪个层次，结论都有实证支撑', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '核心层次：定性空间对应验证', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 四种裂缝宽度梯度(0.1/0.5/1.0/2.0 mm) × 3平行样 = 12块', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 信噪比SNR≥3:1判定"可检出"，目标：>90%岩样可检出荧光 → 定性验证成功', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '探索层次：半定量灰度-裂缝宽度关系', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 裂缝宽度(横坐标) vs 5个ROI平均灰度值(纵坐标)，Spearman秩相关系数ρ', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 若ρ≥0.85且p<0.01 → 灰度可作为裂缝宽度的半定量判据', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 若ρ<0.6 → 荧光信号仅保留"有/无"定性判读功能', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 重要边界条件（防止不当外推）', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '该半定量关系受控于：同一砂岩类型 + 同一成像参数 + 同一荧光粉批次', 'font_size': Pt(15), 'bold': True, 'color': BLACK},
    {'text': '不具备跨条件的绝对定量测量能力', 'font_size': Pt(15), 'bold': True, 'color': ACCENT_GOLD},
]
add_multiline_textbox(slide, 500000, 550000, 8200000, 4400000,
                      quant_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# ═══════════════════════════════
# SLIDE 24: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 6)

# ═══════════════════════════════
# SLIDE 25-26: 现场施工工艺与经济性
# ═══════════════════════════════
# --- Slide 25: 工艺方案 ---
slide = add_blank_slide(prs)
add_section_title(slide, '七、"母液预配+在线稀释"现场施工工艺方案')
add_top_accent_line(slide)
add_bottom_bar(slide)

field_lines = [
    {'text': '▎ 工艺设计思路：避免现场干粉直接加入——粉尘污染+分散不均+结块', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '步骤一：母液预配（配液区）', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 专用配制罐(5~10 m³) + 高速搅拌500~1000 rpm × 30 min + 工业超声40 kHz', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 浓度40 g/L，静置2 h分层≤5%方可放行，有效期7天', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '步骤二：在线稀释（混砂车上游）', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 隔膜计量泵(精度±1%) → SMX型静态混合器(DN50, ≥1 m) → 0.5% v/v在线注入HPG主路', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 混合后进入混砂车汇合支撑剂 → 高压泵组注入井筒', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '步骤三：在线质量控制', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 每30 min取样：便携式荧光分光光度计(520 nm) + 便携式粘度计', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 容许偏差：浓度±15%，粘度±10%空白HPG基液值', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 580000, 8400000, 4300000,
                      field_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 26: 经济性与环保 ---
slide = add_blank_slide(prs)
add_section_title(slide, '七、经济性对比分析与环保安全性')
add_top_accent_line(slide)
add_bottom_bar(slide)

econ_lines = [
    {'text': '▎ 单井材料估算（四川盆地典型页岩气水平井：8段/16000 m³）', 'font_size': Pt(18), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 荧光母液：80 m³（0.5% v/v稀释比）→ 荧光粉3200 kg（40 g/L母液浓度）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 改性试剂：KH550约64 kg + PEG4000约96 kg + 柠檬酸3.2 kg + Triton X-100 1.6 kg', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 经济性对比', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 本方案单井总成本：约100万元（~14万美元）——最大用量方案', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 化学示踪剂参考价：3~8万美元/井（含材料+注入+实验室分析）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 经济有利场景：仅需1~2关键段使用 → 成本降至25~50万元', 'font_size': Pt(15), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '• 间接效益：实物证据验证可避免错误压裂设计决策，节省成本远超示踪方案本身', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 环保安全性', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• Sr²⁺生物毒性类似Ca²⁺，Eu/Dy属轻稀土低毒性(LD50>2000 mg/kg)', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 返排液处理：絮凝-沉降-过滤(13 μm颗粒) + 化学沉淀(调pH 9~10)或离子交换', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '• 含荧光粉岩屑属一般工业固废，SrAl₂O₄基质化学性质稳定、环境浸出风险低', 'font_size': Pt(14), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 适用条件：砂岩储层(>0.1 mm裂缝) + HPG压裂液体系 + 需可实施取心', 'font_size': Pt(16), 'bold': True, 'color': MEDIUM_BLUE},
]
add_multiline_textbox(slide, 500000, 550000, 8200000, 4350000,
                      econ_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(20))

# ═══════════════════════════════
# SLIDE 27: 目录过渡
# ═══════════════════════════════
add_toc_slide(prs, toc_items, 7)

# ═══════════════════════════════
# SLIDE 28-29: 创新点与结论
# ═══════════════════════════════
# --- Slide 28: 创新点 ---
slide = add_blank_slide(prs)
add_section_title(slide, '八、创新点')
add_top_accent_line(slide)
add_bottom_bar(slide)

innov_lines = [
    {'text': '创新点一：面向氧化破胶环境的无机荧光示踪方案', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '选择无机铝酸盐晶体SrAl₂O₄:Eu²⁺,Dy³⁺为示踪介质——发光中心受刚性晶格保护，', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '对过硫酸铵氧化破胶剂具有本征化学惰性——区别于有机荧光树脂和半导体量子点。', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '创新点二："螯合-润湿-位阻"多层次协同分散体系', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '柠檬酸螯合屏蔽Al³⁺（抑制交联干扰）+ 非离子表面活性剂界面润湿（促进PEG伸展）', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '+ PEG长链空间位阻（渗透排斥+构象熵效应）——三机制协同解决高密度荧光粉悬浮问题。', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '创新点三：利用压裂工程时序的功能切换设计', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '提出"破胶诱导PEG脱附→KH550氨基暴露→砂岩硅羟基锚定"的时序功能切换假说，', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '将"注入→关井→返排"工程时间序列转化为颗粒表面化学的功能转化序列。', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '创新点四：动态驱替全流程验证实验方法', 'font_size': Pt(20), 'bold': True, 'color': MEDIUM_BLUE},
    {'text': '构建"注入—关井破胶—返排—取心—成像"五步串联动态实验系统，以驱替压力曲线、', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '返排浓度曲线和岩心成像三组数据构成完整证据链，区别于已有研究的静态评价方法。', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 550000, 8200000, 4350000,
                      innov_lines, Pt(15), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# --- Slide 29: 结论 ---
slide = add_blank_slide(prs)
add_section_title(slide, '八、主要结论与研究展望')
add_top_accent_line(slide)
add_bottom_bar(slide)

conclusion_lines = [
    {'text': '▎ 主要结论（待全部实验完成后撰写）', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '1. 建立了SrAl₂O₄:Eu²⁺,Dy³⁺荧光粉在模拟储层环境中的适用性边界条件，', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '   获取了热稳定性、化学稳定性和水解动力学的关键参数。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '2. 成功开发了"KH550化学锚固+PEG4000物理屏蔽"双层表面改性方案，', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '   建立"柠檬酸+Triton X-100+PEG"多层次协同分散体系，解决了高密度微米粉在HPG基液中的悬浮稳定性问题。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '3. 研制了浓度≥40 g/L的高稳定荧光粉悬浮母液，经SY/T标准验证，', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '   荧光母液的引入不显著损害压裂液的基液、交联、悬砂和破胶性能。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '4. 揭示了破胶诱导PEG脱附-KH550氨基暴露-砂岩硅羟基锚定的功能切换机制。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '5. 通过动态驱替全流程实验验证了荧光压裂液的可注入性、附着牢固度和空间对应性，', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '   在受控条件下建立了裂缝宽度与荧光灰度值的半定量单调映射关系。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '6. 设计了"母液预配+在线稀释"现场施工工艺方案，完成技术经济与环保可行性评估。', 'font_size': Pt(16), 'bold': False, 'color': BLACK},
    {'text': '', 'font_size': Pt(6), 'bold': False, 'color': BLACK},
    {'text': '▎ 展望', 'font_size': Pt(20), 'bold': True, 'color': DARK_BLUE},
    {'text': '• 碳酸盐岩/页岩储层的壁面锚定效果需单独验证', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 井下紫外探头直接观测可作远期替代检测手段', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
    {'text': '• 稀土元素长期环境迁移行为需持续跟踪', 'font_size': Pt(15), 'bold': False, 'color': BLACK},
]
add_multiline_textbox(slide, 500000, 580000, 8200000, 4300000,
                      conclusion_lines, Pt(16), False, BLACK, PP_ALIGN.LEFT, Pt(22))

# ═══════════════════════════════
# SLIDE 30: 致谢
# ═══════════════════════════════
slide = add_blank_slide(prs)
add_bottom_bar(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(9144000), Emu(80000))
shape.fill.solid(); shape.fill.fore_color.rgb = DARK_BLUE; shape.line.fill.background()

add_textbox(slide, 2000000, 1500000, 5000000, 1000000,
            '请各位老师批评指正', Pt(36), True, DARK_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 2000000, 2600000, 5000000, 600000,
            '谢谢！', Pt(28), False, MEDIUM_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 2000000, 3400000, 5000000, 500000,
            '成都理工大学能源学院\n石油与天然气工程专业\n郝乐乐', Pt(16), False, DARK_BLUE, PP_ALIGN.CENTER)

# ── SAVE ──
output_path = r'c:\Users\郝\Desktop\claude\output\答辩PPT_荧光压裂液_郝乐乐.pptx'
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f'✅ PPT saved to: {output_path}')
print(f'   Total slides: {len(prs.slides)}')
