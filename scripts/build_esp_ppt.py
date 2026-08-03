#!/usr/bin/env python3
"""Build ESP-T paper presentation PPTX for advisor reporting.
Follows nature-paper2ppt materials arc (design-to-performance).
No figure images available from DOCX source — use PPT-native tables and structured layouts."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Nature-style restrained academic) ──────────────
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x1A)
BODY_TEXT   = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0x8A)  # muted blue for headings
ACCENT_ORANGE = RGBColor(0xC0, 0x5C, 0x2C) # muted orange for highlights
ACCENT_GREEN  = RGBColor(0x3A, 0x7D, 0x44) # for positive metrics
LIGHT_GRAY    = RGBColor(0x8A, 0x8A, 0x8A)
VERY_LIGHT_BG = RGBColor(0xF5, 0xF3, 0xEF) # warm off-white
TABLE_HEADER_BG = RGBColor(0x2B, 0x5C, 0x8A)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF4)
BORDER_LIGHT  = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ───────────────────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text="",
                font_size=14, bold=False, color=BODY_TEXT,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei",
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    return tf

def add_line(slide, x1, y1, x2, y2, color=BORDER_LIGHT, width=1):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_metric_card(slide, left, top, width, height, value, label, value_color=DARK_TEXT):
    """Add a metric highlight card: big number + small label"""
    add_rect(slide, left, top, width, height, fill_color=VERY_LIGHT_BG)
    tf = add_textbox(slide, left + 0.15, top + 0.15, width - 0.3, height * 0.55,
                     value, font_size=22, bold=True, color=value_color,
                     alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.15, top + height * 0.55, width - 0.3, height * 0.35,
                label, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

def add_slide_number(slide, num, total):
    add_textbox(slide, 12.3, 7.05, 0.8, 0.3, f"{num}/{total}",
                font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_slide_title(slide, title_text, subtitle_text=None):
    """Add a consistent slide title bar at the top"""
    add_rect(slide, 0, 0, 13.333, 0.06, fill_color=ACCENT_BLUE)
    add_textbox(slide, 0.6, 0.25, 12, 0.55, title_text, font_size=26, bold=True, color=DARK_TEXT)
    if subtitle_text:
        add_textbox(slide, 0.6, 0.78, 12, 0.35, subtitle_text, font_size=12, color=LIGHT_GRAY)
    # separator line
    add_line(slide, 0.6, 1.15, 12.7, 1.15, color=BORDER_LIGHT, width=0.5)

TOTAL_SLIDES = 14

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
# Dark accent band at top
add_rect(slide, 0, 0, 13.333, 0.08, fill_color=ACCENT_BLUE)
# Left accent strip
add_rect(slide, 0.8, 1.5, 0.06, 4.5, fill_color=ACCENT_ORANGE)

# Title
add_textbox(slide, 1.3, 1.6, 11, 1.4,
            "环氧树脂微球包覆亲油性 Fe₃O₄ 纳米颗粒\n示踪支撑剂用于分段产量监测",
            font_size=30, bold=True, color=DARK_TEXT)

# English title
add_textbox(slide, 1.3, 3.2, 11, 0.8,
            "Epoxy Resin Microspheres Encapsulating Oleophilic Fe₃O₄\nNanoparticles as Tracer Proppants for Production Allocation",
            font_size=14, color=LIGHT_GRAY)

# Separator
add_line(slide, 1.3, 4.2, 5.5, 4.2, color=ACCENT_ORANGE, width=2)

# Meta info
add_textbox(slide, 1.3, 4.5, 8, 0.4, "ESP-T: Epoxy resin microSpheres with oleophilic Tracer",
            font_size=13, color=ACCENT_BLUE)
add_textbox(slide, 1.3, 5.1, 8, 0.4, "导师汇报  |  2026年6月",
            font_size=12, color=LIGHT_GRAY)
add_textbox(slide, 1.3, 5.5, 8, 0.4, "材料 / 石油工程 — 非常规储层压裂示踪技术",
            font_size=11, color=LIGHT_GRAY)

# Bottom accent
add_rect(slide, 0, 7.42, 13.333, 0.08, fill_color=ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: BACKGROUND — Why does this matter?
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "分段产量监测是压裂效果评价的核心需求")
add_slide_number(slide, 2, TOTAL_SLIDES)

# Left: Problem statement
add_textbox(slide, 0.6, 1.5, 5.5, 0.4, "核心矛盾", font_size=16, bold=True, color=ACCENT_ORANGE)
tf = add_textbox(slide, 0.6, 2.0, 5.5, 2.0,
                 "• 压裂后各段产油贡献差异显著\n"
                 "• 各段产量以不同速率衰减\n"
                 "• 缺乏精确的分段产量分配手段\n"
                 "• 导致开发策略优化缺乏数据支撑",
                 font_size=13, color=BODY_TEXT)
tf.paragraphs[0].space_after = Pt(8)

# Right: Key fact cards
add_metric_card(slide, 7.2, 1.6, 2.5, 1.0,
                ">50%", "非常规油气占全球\n剩余可采储量比例", ACCENT_BLUE)
add_metric_card(slide, 10.0, 1.6, 2.7, 1.0,
                "357.27°C", "ESP-T 初始分解温度\n远超井下需求 (80-200°C)", ACCENT_BLUE)

# Bottom: Current technology limitations
add_textbox(slide, 0.6, 4.4, 12, 0.4, "现有示踪技术的三大瓶颈", font_size=16, bold=True, color=ACCENT_ORANGE)

# Three limitation cards
lims = [
    ("传统油溶性示踪剂", "与压裂液兼容性差\n无法实现长期监测"),
    ("涂层型示踪支撑剂", "涂层溶解即终止监测\n多步骤制备工艺复杂"),
    ("PS微球示踪剂", "机械强度差\n热稳定性不足 (<200°C)"),
]
for i, (title, desc) in enumerate(lims):
    x = 0.6 + i * 4.2
    add_rect(slide, x, 5.0, 3.8, 1.8, fill_color=VERY_LIGHT_BG)
    add_textbox(slide, x + 0.2, 5.1, 3.4, 0.4, title, font_size=13, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, x + 0.2, 5.55, 3.4, 1.0, desc, font_size=11, color=BODY_TEXT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: GAP — The research opportunity
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "环氧树脂作为油相示踪释放基体的研究空白")
add_slide_number(slide, 3, TOTAL_SLIDES)

# Literature landscape
add_textbox(slide, 0.6, 1.5, 12, 0.35, "已有研究 vs 本文切入点", font_size=16, bold=True, color=ACCENT_ORANGE)

# Table: literature landscape
rows_data = [
    ["研究", "基体材料", "示踪剂类型", "监测相态", "局限性"],
    ["Li et al. 2021 [26]", "环氧树脂", "水溶性示踪剂", "💧 水相", "仅水相释放"],
    ["Wei et al. 2024 [27]", "环氧树脂", "水溶性示踪剂", "💧 水相", "定位产水层段"],
    ["Gong et al. 2024 [19]", "聚苯乙烯(PS)", "油溶性示踪剂", "🛢️ 油相", "PS强度/热稳定性不足"],
    ["本文 ESP-T", "环氧树脂", "亲油 nano-Fe₃O₄@SA", "🛢️ 油相 ✓", "首次油相监测 ✓"],
]
table = slide.shapes.add_table(len(rows_data), 5, Inches(0.6), Inches(2.1), Inches(12.1), Inches(2.6)).table
for r, row in enumerate(rows_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r == len(rows_data) - 1:
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            else:
                p.font.color.rgb = BODY_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r == len(rows_data) - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = VERY_LIGHT_BG

# Key insight
add_textbox(slide, 0.6, 5.1, 12, 0.6,
            "▸ 核心创新：首次以环氧树脂为亲油性释放基体，通过硬脂酸改性实现 nano-Fe₃O₄@SA 的均匀包覆",
            font_size=14, bold=True, color=ACCENT_ORANGE)

# Innovation highlights
add_textbox(slide, 0.6, 5.8, 12, 0.4, "三合一思路", font_size=15, bold=True, color=ACCENT_BLUE)
items = [
    "✅  乳液聚合一步合成 —— 低密度调控 + 纳米改性 + 原位包覆",
    "✅  硬脂酸表面改性 —— 亲油化，WCA 从 72.3° → 104.6°",
    "✅  ADE分段模型 —— 物理意义清晰，R² = 0.9939",
]
for i, item in enumerate(items):
    add_textbox(slide, 0.8 + (i % 3) * 4.1, 6.3, 3.8, 0.4, item, font_size=11, color=BODY_TEXT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: EXPERIMENTAL DESIGN — Synthesis & Characterization
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "ESP-T 的合成路线与表征体系")
add_slide_number(slide, 4, TOTAL_SLIDES)

# Synthesis flow — horizontal process
add_textbox(slide, 0.6, 1.5, 5, 0.35, "▎合成路线", font_size=15, bold=True, color=ACCENT_BLUE)

steps = [
    ("① 共沉淀", "FeCl₃ + FeCl₂\n+ Mn²⁺掺杂\nNH₃·H₂O, pH=10\n80°C, 2h"),
    ("② 表面改性", "硬脂酸乙醇溶液\n超声亲油改性\n→ nano-Fe₃O₄@SA"),
    ("③ 乳液聚合", "E51 + T31 + 中空玻璃微球\nSiO₂分散液 + 胍胶\n380 RPM 搅拌"),
    ("④ 固化干燥", "50°C 固化 1h\n去离子水冲洗\n80°C 干燥 10h"),
]
for i, (title, desc) in enumerate(steps):
    x = 0.6 + i * 3.2
    add_rect(slide, x, 2.0, 2.8, 2.2, fill_color=VERY_LIGHT_BG)
    add_textbox(slide, x + 0.15, 2.1, 2.5, 0.35, title, font_size=13, bold=True, color=ACCENT_ORANGE)
    add_textbox(slide, x + 0.15, 2.55, 2.5, 1.4, desc, font_size=10, color=BODY_TEXT)
    if i < 3:
        add_textbox(slide, x + 2.85, 2.7, 0.3, 0.5, "→", font_size=20, bold=True, color=ACCENT_BLUE)

# Characterization methods
add_textbox(slide, 0.6, 4.6, 5, 0.35, "▎表征与测试体系", font_size=15, bold=True, color=ACCENT_BLUE)

chars = [
    ("SEM + EDS", "形貌 / 元素分布"),
    ("TGA / DSC", "热稳定性 / 分解行为"),
    ("WCA", "润湿性 (水接触角)"),
    ("物理力学", "密度 / 圆球度 / 破碎率 / 酸溶"),
    ("油水过滤", "导流能力间接评价"),
    ("ICP-MS", "示踪释放浓度定量"),
    ("K-P 模型", "释放动力学机制"),
    ("ADE 模型", "突破曲线拟合 / 产量定量"),
]
for i, (method, purpose) in enumerate(chars):
    x = 0.6 + (i % 4) * 3.2
    y = 5.1 + (i // 4) * 0.9
    add_rect(slide, x, y, 2.9, 0.75, fill_color=VERY_LIGHT_BG)
    add_textbox(slide, x + 0.1, y + 0.05, 1.3, 0.3, method, font_size=11, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, x + 1.4, y + 0.05, 1.4, 0.3, purpose, font_size=9, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: SEM RESULTS — Morphology
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "SEM 表征：nano-Fe₃O₄@SA 纳米团簇均匀分散于环氧基体中")
add_slide_number(slide, 5, TOTAL_SLIDES)

# Side-by-side comparison
add_textbox(slide, 0.6, 1.5, 5.5, 0.4, "纯环氧微球", font_size=15, bold=True, color=LIGHT_GRAY)
add_textbox(slide, 7.0, 1.5, 5.5, 0.4, "ESP-T（nano-Fe₃O₄@SA 掺杂）", font_size=15, bold=True, color=ACCENT_ORANGE)

# Placeholder boxes for SEM images
add_rect(slide, 0.6, 2.0, 5.5, 3.5, fill_color=VERY_LIGHT_BG, line_color=BORDER_LIGHT)
add_textbox(slide, 1.5, 3.4, 3.7, 0.5, "[SEM 图片占位 — Fig. 3-1 上层]\n表面光滑，少量褶皱", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, 7.0, 2.0, 5.5, 3.5, fill_color=VERY_LIGHT_BG, line_color=ACCENT_ORANGE)
add_textbox(slide, 8.0, 3.4, 3.5, 0.5, "[SEM 图片占位 — Fig. 3-1 下层]\n岛状凸起，纳米团簇均匀分布", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key observations
add_textbox(slide, 0.6, 5.8, 12, 0.3, "▸ 关键发现", font_size=14, bold=True, color=ACCENT_BLUE)
findings = [
    "❶  纳米团簇以嵌入键合态存在，无界面裂缝 —— 硬脂酸长烷基链与环氧分子链形成物理缠结",
    "❷  固化收缩过程中纳米团簇被挤出至表面 → 表面富集的粗糙纳米结构",
    "❸  Fe 元素遍布整个颗粒 (EDS Mapping) → 封装成功的直接证据",
]
for i, f in enumerate(findings):
    add_textbox(slide, 0.6 + (i % 3) * 4.1, 6.2, 3.9, 0.6, f, font_size=10, color=BODY_TEXT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: PERFORMANCE DASHBOARD — Thermal, Wettability, Density
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "ESP-T 性能仪表盘：耐高温、疏水、低密度、高强度")
add_slide_number(slide, 6, TOTAL_SLIDES)

# Top row: 3 big metric cards
add_metric_card(slide, 0.6, 1.5, 3.8, 1.8,
                "357.27 °C", "初始分解温度\n远超井下工况 (80-200°C)\n仅微量水分挥发 <200°C", ACCENT_BLUE)
add_metric_card(slide, 4.8, 1.5, 3.8, 1.8,
                "104.6°", "水接触角 (WCA)\n从 72.3° 跃升 32.3°\n亲水 → 疏水转变", ACCENT_ORANGE)
add_metric_card(slide, 9.0, 1.5, 3.8, 1.8,
                "0.646 g/cm³", "体积密度\n＜水 (1 g/cm³)\n压裂液中可悬浮", ACCENT_GREEN)

# Middle row: performance table
add_textbox(slide, 0.6, 3.6, 5, 0.35, "▎行业标准达标情况", font_size=14, bold=True, color=ACCENT_BLUE)

perf_data = [
    ["指标", "ESP-T 实测值", "行业标准", "状态"],
    ["球度 / 圆度", "> 0.9", "Krumbien-Sloss", "✅ 达标"],
    ["酸溶解度", "3.3%", "≤ 5% (SY/T 5107)", "✅ 达标"],
    ["52 MPa 破碎率", "2.9%", "—", "✅ 与纯环氧 (2.6%) 相当"],
    ["环压焓 ΔH", "99.53 J/g", "—", "定量表征"],
]
table = slide.shapes.add_table(len(perf_data), 4, Inches(0.6), Inches(4.1), Inches(8.5), Inches(2.0)).table
for r, row in enumerate(perf_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = BODY_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT

# Bottom: TGA detail
add_textbox(slide, 9.8, 4.1, 3.0, 0.3, "▎TGA 三阶段分解", font_size=12, bold=True, color=ACCENT_BLUE)
tga_stages = [
    "I:   50-350°C  失重 5.70% (水分/乙醇)",
    "II:  350-400°C 失重 72.5% (主分解)",
    "III: >400°C  残渣: 玻球+Fe₃O₄@SA",
]
for i, stage in enumerate(tga_stages):
    add_textbox(slide, 9.8, 4.5 + i * 0.35, 3.0, 0.3, stage, font_size=9, color=BODY_TEXT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: WATER-RESISTANT, OIL-PERMEABLE — The key property
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "阻水亲油特性：油过滤时间缩短 66%，水过滤时间延长 10 倍")
add_slide_number(slide, 7, TOTAL_SLIDES)

# Comparison table — the visual centerpiece
comp_data = [
    ["", "纯环氧微球", "ESP-T", "变化"],
    ["水接触角 (WCA)", "72.3° (弱亲水)", "104.6° (疏水)", "↑ 32.3°"],
    ["水过滤时间", "2 分 53 秒", "28 分 41 秒", "↑ 10 倍"],
    ["油过滤时间", "15 分 11 秒", "5 分 11 秒", "↓ 66.1%"],
]
table = slide.shapes.add_table(len(comp_data), 4, Inches(0.6), Inches(1.5), Inches(7.5), Inches(2.0)).table
for r, row in enumerate(comp_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13 if r > 0 else 12)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif c == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE if "↑" in cell_text else ACCENT_GREEN
            else:
                p.font.color.rgb = BODY_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG

# Mechanism diagram (text)
add_textbox(slide, 9.0, 1.5, 3.8, 0.35, "▎机理", font_size=14, bold=True, color=ACCENT_BLUE)
add_rect(slide, 9.0, 2.0, 3.8, 4.5, fill_color=VERY_LIGHT_BG)

mech_text = (
    "硬脂酸 C₁₇H₃₅COOH\n\n"
    "• -COOH → 与 Fe₃O₄ 表面\n"
    "  羟基配位成键\n\n"
    "• -C₁₇H₃₅ 长烷基链 →\n"
    "  向外定向排列\n\n"
    "→ 构建疏水膜\n"
    "→ 降低表面自由能\n\n"
    "水相：被排斥 → 阻力↑\n"
    "油相：相容铺展 → 阻力↓"
)
add_textbox(slide, 9.2, 2.1, 3.4, 4.2, mech_text, font_size=11, color=BODY_TEXT)

# Bottom takeaway
add_textbox(slide, 0.6, 5.8, 12, 0.4,
            "▸ 工程意义：支撑剂充填层为油提供增强导流 → 促进油流向井筒；水相阻力增大 → 缓解水窜，提高采收率",
            font_size=13, bold=True, color=ACCENT_ORANGE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: RELEASE KINETICS — Korsmeyer-Peppas model
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "示踪释放动力学：Fick 扩散与 Case-II 松弛协同主导")
add_slide_number(slide, 8, TOTAL_SLIDES)

# K-P equation
add_textbox(slide, 0.6, 1.5, 6, 0.35, "Korsmeyer–Peppas 模型", font_size=15, bold=True, color=ACCENT_BLUE)
add_textbox(slide, 0.6, 1.9, 3, 0.8, "C/C₀ = K · tⁿ", font_size=22, bold=True, color=DARK_TEXT)
add_textbox(slide, 3.6, 1.95, 7, 0.6,
            "n ≤ 0.43 → Fick 扩散\n"
            "0.43 < n < 0.85 → 非 Fick 异常传输 ⬅ 本文\n"
            "n ≥ 0.85 → Case-II 松弛",
            font_size=11, color=BODY_TEXT)

# K-P fitting table
kp_data = [
    ["参数", "30°C", "60°C", "90°C", "120°C", "趋势"],
    ["R²", "0.9549", "0.9649", "0.9560", "0.9454", "> 0.94"],
    ["K", "0.0554", "0.0818", "0.1134", "0.1964", "温度↑ → K↑"],
    ["n", "0.5983", "0.6665", "0.5684", "0.5557", "0.45–0.85"],
]
table = slide.shapes.add_table(len(kp_data), 6, Inches(0.6), Inches(2.9), Inches(12.1), Inches(1.8)).table
for r, row in enumerate(kp_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = BODY_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r == 2:  # K row — highlight
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT

# Release mechanism interpretation
add_textbox(slide, 0.6, 5.0, 12, 0.35, "▸ 释放机制推断", font_size=14, bold=True, color=ACCENT_BLUE)
mech_items = [
    "十二烷渗透进入交联环氧网络 → 溶胀产生玻璃核 + 凝胶层双层结构",
    "溶胀降低聚合物链缠结 → 形成扩展传输通道 → 示踪剂扩散至外部介质",
    "升温增强溶剂渗透速率 → 加速溶胀 → 增大孔径 → 促进扩散",
]
for i, item in enumerate(mech_items):
    add_textbox(slide, 0.6 + (i % 3) * 4.1, 5.5, 3.9, 0.7, item, font_size=10, color=BODY_TEXT)

add_textbox(slide, 0.6, 6.2, 12, 0.35,
            "⚠ K-P 模型为幂律形式，严格适用于释放初期 (Mₜ/M∞ < 0.6)，不应外推至实验时限外",
            font_size=9, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: ADE MODEL — The crown jewel
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "ADE 分段模型：R² = 0.9939，erfc 拖尾贡献 47% 信号")
add_slide_number(slide, 9, TOTAL_SLIDES)

# Model structure
add_textbox(slide, 0.6, 1.5, 7, 0.35, "分段对流-弥散模型 (tanh 平滑过渡)", font_size=15, bold=True, color=ACCENT_BLUE)
add_textbox(slide, 0.6, 1.9, 12, 0.5,
            "C(t) = cb + A·C_rise(t) + a·C_fall(t)     (tanh 权重函数平滑拼接)",
            font_size=16, bold=True, color=DARK_TEXT)
add_textbox(slide, 0.6, 2.35, 12, 0.3,
            "C_rise: ADE 瞬时脉冲解 (Gaussian 形式) — 关井积蓄示踪剂团  |  C_fall: ADE 连续源解 (erfc 拖尾) — 基质扩散控制持续释放",
            font_size=11, color=LIGHT_GRAY)

# Key results table (the centerpiece)
result_data = [
    ["参数", "拟合值", "参考值/标准", "验证"],
    ["R²", "0.9939", "—", "—"],
    ["RMSE", "0.0210", "—", "残差 < ±2σ"],
    ["拟合流量 Q", "0.46 mL/min", "泵设 0.50 mL/min", "误差 8% ✓"],
    ["平均停留时间 MRT", "37.4 min", "对流时间 x/v = 38.6 min", "比值 0.967 ✓"],
    ["Peclet 数", "0.934 (≈1)", "Pe > 1 = 对流主导", "过渡态，与缓释一致 ✓"],
    ["erfc 拖尾占比", "47%", "—", "长期监测的物理基础"],
]
table = slide.shapes.add_table(len(result_data), 4, Inches(0.6), Inches(2.8), Inches(9.5), Inches(2.6)).table
for r, row in enumerate(result_data):
    for c, cell_text in enumerate(row):
        cell = table.cell(r, c)
        cell.text = cell_text
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Microsoft YaHei"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r == len(result_data) - 1:
                p.font.bold = True
                p.font.color.rgb = ACCENT_ORANGE
            else:
                p.font.color.rgb = BODY_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif r == len(result_data) - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = VERY_LIGHT_BG

# Right: 47% highlight card
add_rect(slide, 10.5, 2.8, 2.3, 2.6, fill_color=VERY_LIGHT_BG)
add_textbox(slide, 10.7, 3.1, 1.9, 0.8, "47%", font_size=36, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 10.7, 4.0, 1.9, 0.8, "erfc 拖尾\n信号占比", font_size=11, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 10.7, 4.7, 1.9, 0.5, "≈ 一半信号来自\n基质缓释 🎯", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom takeaway
add_textbox(slide, 0.6, 5.7, 12, 0.5,
            "▸ 47% 的意义：基质扩散控制释放是长期监测的物理基础 —— 即使在流动条件下，非 Fick 机制仍主导示踪传输",
            font_size=13, bold=True, color=ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: TWO-PHASE MONITORING
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "两相流验证：示踪剂通量 FO 定量追踪油相产量")
add_slide_number(slide, 10, TOTAL_SLIDES)

# Experimental conditions
add_textbox(slide, 0.6, 1.5, 7, 0.35, "实验条件", font_size=14, bold=True, color=ACCENT_BLUE)
add_textbox(slide, 0.6, 1.9, 12, 0.5,
            "• 油水比 (OWR): 4:1 / 1:1 / 1:4    • 总两相流量: 0.1 / 0.2 / 0.3 / 0.4 mL/min    • 稳态连续生产（无闷井）",
            font_size=12, color=BODY_TEXT)

# Key findings — two-column
add_textbox(slide, 0.6, 2.7, 5.5, 0.3, "示踪剂浓度 vs 总流量 & OWR", font_size=13, bold=True, color=ACCENT_ORANGE)
add_textbox(slide, 0.6, 3.05, 5.5, 1.5,
            "• 浓度随总流量增大而降低\n"
            "  → 稀释效应：单位体积接触时间缩短\n\n"
            "• 浓度与 OWR 关系不大\n"
            "  → 但 FO 随 OWR 增大而升高\n\n"
            "• FO 与总流量无关\n"
            "  → 恒定 OWR 下接触面积不变",
            font_size=12, color=BODY_TEXT)

add_textbox(slide, 7.0, 2.7, 5.5, 0.3, "归一化 FO ≈ 实际油相流量", font_size=13, bold=True, color=ACCENT_ORANGE)
add_textbox(slide, 7.0, 3.05, 5.5, 1.5,
            "• 归一化 FO 以单相油驱稳态 FO\n"
            "  (3.187 μg/min) 为基准标定\n\n"
            "• 各 OWR 下 FO 与实际油相流量\n"
            "  高度吻合 ✓\n\n"
            "→ 恒定总流量下可从 FO 变化曲线\n"
            "   量化各段产油量",
            font_size=12, color=BODY_TEXT)

# Bottom: takeaway with FO definition
add_rect(slide, 0.6, 5.0, 12.1, 0.7, fill_color=VERY_LIGHT_BG)
add_textbox(slide, 0.8, 5.1, 11.7, 0.5,
            "FO = 示踪剂通量 = 单位时间通过井口采样点的示踪剂质量 = ESP-T 示踪释放速率",
            font_size=12, bold=True, color=DARK_TEXT)

add_textbox(slide, 0.6, 5.9, 12, 0.5,
            "▸ 在稳态两相流条件下，可通过 FO 变化曲线量化标记层段的油相流量 → 为分段产量分配提供实践基础",
            font_size=12, color=ACCENT_BLUE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: STRUCTURE-PROPERTY RELATIONSHIP
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "构效关系：从微观结构到宏观性能的完整链条")
add_slide_number(slide, 11, TOTAL_SLIDES)

# Central chain
chain_items = [
    ("硬脂酸改性\nnano-Fe₃O₄@SA", "ACCENT_ORANGE"),
    ("长烷基链\n物理缠结", "ACCENT_BLUE"),
    ("纳米团簇\n嵌入键合", "ACCENT_BLUE"),
    ("疏水表面\nWCA 104.6°", "ACCENT_ORANGE"),
    ("阻水亲油\n导流特性", "ACCENT_GREEN"),
    ("非Fick释放\nn=0.45–0.85", "ACCENT_ORANGE"),
    ("ADE模型\nR²=0.9939", "ACCENT_BLUE"),
    ("分段产量\n监测", "ACCENT_GREEN"),
]
colors_map = {"ACCENT_ORANGE": ACCENT_ORANGE, "ACCENT_BLUE": ACCENT_BLUE, "ACCENT_GREEN": ACCENT_GREEN}

for i, (text, color_key) in enumerate(chain_items):
    x = 0.3 + i * 1.65
    y = 1.8
    color = colors_map[color_key]
    add_rect(slide, x, y, 1.4, 1.2, fill_color=VERY_LIGHT_BG)
    # top accent line
    add_rect(slide, x, y, 1.4, 0.04, fill_color=color)
    add_textbox(slide, x + 0.05, y + 0.15, 1.3, 0.9, text, font_size=10, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    if i < len(chain_items) - 1:
        add_textbox(slide, x + 1.4, y + 0.3, 0.25, 0.5, "▸", font_size=16, bold=True, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: two key mechanisms
add_textbox(slide, 0.6, 3.5, 12, 0.35, "▎两个核心机制", font_size=14, bold=True, color=ACCENT_BLUE)

add_rect(slide, 0.6, 4.0, 5.8, 2.6, fill_color=VERY_LIGHT_BG)
add_textbox(slide, 0.8, 4.1, 5.4, 0.3, "润湿性转变机制", font_size=13, bold=True, color=ACCENT_ORANGE)
add_textbox(slide, 0.8, 4.5, 5.4, 2.0,
            "-COOH 配位键合 + -C₁₇H₃₅ 向外定向\n"
            "→ 疏水膜构建 → 表面自由能 ↓\n"
            "→ WCA: 72.3° → 104.6°\n"
            "→ 油通过时间 ↓66.1% | 水↑10×\n"
            "→ 阻水亲油 → 选择性导流",
            font_size=11, color=BODY_TEXT)

add_rect(slide, 6.9, 4.0, 5.8, 2.6, fill_color=VERY_LIGHT_BG)
add_textbox(slide, 7.1, 4.1, 5.4, 0.3, "非Fick释放机制", font_size=13, bold=True, color=ACCENT_ORANGE)
add_textbox(slide, 7.1, 4.5, 5.4, 2.0,
            "溶剂渗透 → 环氧网络溶胀\n"
            "→ 玻璃核 + 凝胶层双层结构\n"
            "→ 聚合物链缠结 ↓ → 传输通道 ↑\n"
            "→ n = 0.45–0.85: Fick + Case-II 协同\n"
            "→ 升温 → K↑ → 加速释放",
            font_size=11, color=BODY_TEXT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: LIMITATIONS
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "局限性与当前适用边界")
add_slide_number(slide, 12, TOTAL_SLIDES)

lims = [
    ("单段假设", "ADE 模型假设单裂缝段 + 均匀填充\n多段相互作用 & 非均质性 → 矿场应用偏差"),
    ("稳态假设", "示踪剂通量标定依赖稳态流动\n瞬变流态 (开/关井, 快速降压) 下 FO-Oil 关系可能不成立"),
    ("模型流体", "实验室用十二烷模拟原油\n沥青质吸附、粘度变化等原油组分效应未评估"),
    ("工况边界", ">120°C 高温 + 高矿化度 + CO₂/H₂S\n环氧基体长期化学稳定性 & 硬脂酸改性完整性待验证"),
]
for i, (title, desc) in enumerate(lims):
    x = 0.5 + (i % 2) * 6.2
    y = 1.5 + (i // 2) * 2.7
    add_rect(slide, x, y, 5.9, 2.3, fill_color=VERY_LIGHT_BG)
    # Side accent
    add_rect(slide, x, y, 0.05, 2.3, fill_color=ACCENT_ORANGE)
    add_textbox(slide, x + 0.3, y + 0.2, 5.3, 0.35, f"⚠ {title}", font_size=15, bold=True, color=ACCENT_ORANGE)
    add_textbox(slide, x + 0.3, y + 0.7, 5.3, 1.3, desc, font_size=11, color=BODY_TEXT)

add_textbox(slide, 0.6, 6.5, 12, 0.5,
            "▸ 这些局限性界定了 ESP-T 当前的适用性范围，也为后续研究指明了方向",
            font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13: CONCLUSIONS
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_title(slide, "总结：ESP-T 实现支撑 + 监测双功能一体化")
add_slide_number(slide, 13, TOTAL_SLIDES)

conclusions = [
    ("❶ 材料合成", "nano-Fe₃O₄@SA 均匀包覆于环氧基体\n球度>0.9，满足工业标准"),
    ("❷ 性能突破", "耐温 357°C | 疏水 WCA 104.6°\n阻水亲油 | 密度<水 | 破碎率 2.9%"),
    ("❸ 释放机制", "K-P 模型 R²>0.94 | n=0.45–0.85\nFick + Case-II 协同非Fick传输"),
    ("❹ 监测模型", "ADE tanh 分段模型 R²=0.9939\nerfc 拖尾 47% → 长期监测物理基础"),
    ("❺ 产量定量", "FO 稳态标定 → 油相产量定量\n分段产量分配的实践基础"),
    ("❻ 应用场景", "酸化压裂 / 深井 / 高压\n非常规储层增产 & 长期监测"),
]
for i, (title, desc) in enumerate(conclusions):
    x = 0.4 + (i % 3) * 4.2
    y = 1.5 + (i // 3) * 2.8
    add_rect(slide, x, y, 3.9, 2.5, fill_color=VERY_LIGHT_BG)
    add_rect(slide, x, y, 3.9, 0.04, fill_color=ACCENT_BLUE)
    add_textbox(slide, x + 0.15, y + 0.2, 3.6, 0.35, title, font_size=14, bold=True, color=ACCENT_BLUE)
    add_textbox(slide, x + 0.15, y + 0.7, 3.6, 1.5, desc, font_size=11, color=BODY_TEXT)

# Bottom emphasis
add_rect(slide, 0.6, 6.7, 12.1, 0.6, fill_color=ACCENT_BLUE)
add_textbox(slide, 0.8, 6.75, 11.7, 0.5,
            "ESP-T = 裂缝支撑 + 长期产量监测 → 压裂增产 & 储层管理的双功能平台",
            font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14: THANK YOU
# ═══════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 0.08, fill_color=ACCENT_BLUE)
add_rect(slide, 0.8, 3.0, 0.06, 1.8, fill_color=ACCENT_ORANGE)

add_textbox(slide, 1.3, 3.1, 11, 0.8, "谢谢！请老师批评指正", font_size=32, bold=True, color=DARK_TEXT)
add_textbox(slide, 1.3, 4.0, 11, 0.5,
            "ESP-T: Epoxy Resin Microspheres Encapsulating Oleophilic Fe₃O₄ Nanoparticles\nas Tracer Proppants for Production Allocation",
            font_size=12, color=LIGHT_GRAY)

add_line(slide, 1.3, 5.0, 4.5, 5.0, color=ACCENT_ORANGE, width=2)
add_textbox(slide, 1.3, 5.3, 11, 0.5, "论文全文 & 详细双语对照见: 四氧化三铁环氧树脂拟合/paper.md", font_size=11, color=LIGHT_GRAY)

add_rect(slide, 0, 7.42, 13.333, 0.08, fill_color=ACCENT_BLUE)

# ── SAVE ───────────────────────────────────────────────────────────
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "..", "四氧化三铁环氧树脂拟合", "ESP-T_汇报PPT.pptx")
# Resolve
output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "四氧化三铁环氧树脂拟合", "ESP-T_汇报PPT.pptx"))
os.makedirs(os.path.dirname(output_path), exist_ok=True)

prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")