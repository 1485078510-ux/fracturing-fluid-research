# -*- coding: utf-8 -*-
"""Office file handler -- read/create DOCX, XLSX, PDF, PPTX from CLI.

Usage:
  python office.py read  <file>           输出文件文本内容
  python office.py meta  <file>           输出文件元信息（页数、行列数、幻灯片数等）
  python office.py new   <type> <output>  创建空文件 (docx/xlsx/pdf/pptx)

Examples:
  python office.py read  report.docx
  python office.py read  data.xlsx
  python office.py read  slides.pptx
  python office.py read  paper.pdf --pages 1-5
  python office.py meta  report.docx
"""

import sys, os, json, argparse, io
# Force UTF-8 output on Windows
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

# ── DOCX ────────────────────────────────────────────────────
def read_docx(path):
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))
    return "\n".join(lines)


def meta_docx(path):
    from docx import Document
    doc = Document(path)
    paras = len(doc.paragraphs)
    tables = len(doc.tables)
    sections = len(doc.sections)
    return {
        "type": "docx",
        "paragraphs": paras,
        "tables": tables,
        "sections": sections,
        "size_kb": round(os.path.getsize(path) / 1024, 1),
    }


# ── XLSX ────────────────────────────────────────────────────
def read_xlsx(path, sheet=None):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = [sheet] if sheet else wb.sheetnames
    output = []
    for name in sheets:
        ws = wb[name]
        output.append(f"=== Sheet: {name} ===")
        for row in ws.iter_rows(values_only=True):
            row_str = "\t".join(str(c) if c is not None else "" for c in row)
            if row_str.strip():
                output.append(row_str)
    return "\n".join(output)


def meta_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        sheets[name] = {
            "rows": ws.max_row,
            "cols": ws.max_column,
        }
    return {
        "type": "xlsx",
        "sheets": sheets,
        "size_kb": round(os.path.getsize(path) / 1024, 1),
    }


# ── PDF ─────────────────────────────────────────────────────
def read_pdf(path, pages=None):
    import fitz  # pymupdf
    doc = fitz.open(path)
    output = []
    page_range = parse_page_range(pages, len(doc))

    for i in page_range:
        page = doc[i]  # 0-indexed
        text = page.get_text()
        if text.strip():
            output.append(f"--- Page {i + 1} ---")
            output.append(text)
    return "\n".join(output)


def meta_pdf(path):
    import fitz
    doc = fitz.open(path)
    return {
        "type": "pdf",
        "pages": len(doc),
        "size_kb": round(os.path.getsize(path) / 1024, 1),
        "metadata": dict(doc.metadata),
    }


def parse_page_range(spec, total):
    """Parse page range like '1-5' or '3' returning 0-indexed list."""
    if not spec:
        return list(range(total))
    pages = []
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            pages.extend(range(int(a) - 1, int(b)))
        else:
            pages.append(int(part) - 1)
    return [p for p in pages if 0 <= p < total]


# ── PPTX ────────────────────────────────────────────────────
def read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    output = []
    for i, slide in enumerate(prs.slides, 1):
        output.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        output.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    output.append(" | ".join(cells))
    return "\n".join(output)


def meta_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    return {
        "type": "pptx",
        "slides": len(prs.slides),
        "size_kb": round(os.path.getsize(path) / 1024, 1),
    }


# ── Router ──────────────────────────────────────────────────
READERS = {
    ".docx": read_docx,
    ".xlsx": read_xlsx,
    ".pdf": read_pdf,
    ".pptx": read_pptx,
}

META = {
    ".docx": meta_docx,
    ".xlsx": meta_xlsx,
    ".pdf": meta_pdf,
    ".pptx": meta_pptx,
}


def main():
    parser = argparse.ArgumentParser(description="Office file handler")
    parser.add_argument("action", choices=["read", "meta", "new"])
    parser.add_argument("target", help="file path or output type (for 'new')")
    parser.add_argument("extra", nargs="?", help="output path (for 'new')")
    parser.add_argument("--pages", help="page range for PDF, e.g. 1-5")
    parser.add_argument("--sheet", help="sheet name for XLSX")

    args = parser.parse_args()

    if args.action == "read":
        ext = os.path.splitext(args.target)[1].lower()
        if ext not in READERS:
            print(f"Unsupported format: {ext}", file=sys.stderr)
            sys.exit(1)
        kwargs = {}
        if ext == ".pdf" and args.pages:
            kwargs["pages"] = args.pages
        if ext == ".xlsx" and args.sheet:
            kwargs["sheet"] = args.sheet
        print(READERS[ext](args.target, **kwargs))

    elif args.action == "meta":
        ext = os.path.splitext(args.target)[1].lower()
        if ext not in META:
            print(f"Unsupported format: {ext}", file=sys.stderr)
            sys.exit(1)
        print(json.dumps(META[ext](args.target), ensure_ascii=False, indent=2))

    elif args.action == "new":
        fmt = args.target.lower()
        path = args.extra or f"new_file.{fmt}"
        if fmt == "docx":
            from docx import Document
            Document().save(path)
        elif fmt == "xlsx":
            import openpyxl
            openpyxl.Workbook().save(path)
        elif fmt == "pptx":
            from pptx import Presentation
            Presentation().save(path)
        elif fmt == "pdf":
            import fitz
            fitz.open().save(path)
        else:
            print(f"Unknown type: {fmt}", file=sys.stderr)
            sys.exit(1)
        print(f"Created: {os.path.abspath(path)}")


if __name__ == "__main__":
    main()