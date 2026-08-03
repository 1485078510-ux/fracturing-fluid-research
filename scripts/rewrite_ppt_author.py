#!/usr/bin/env python3
"""Rewrite ALL speaker notes from first-person (author) perspective
for advisor pre-submission meeting. Also adds a journal submission slide."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

pptx_path = os.path.abspath(os.path.join(
    os.path.dirname(__file__), "..", "四氧化三铁环氧树脂拟合", "ESP-T_汇报PPT.pptx"))
prs = Presentation(pptx_path)

# ── Color constants (matching build script) ────────────────────────
ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0x8A)
ACCENT_ORANGE = RGBColor(0xC0, 0x5C, 0x2C)
ACCENT_GREEN = RGBColor(0x3A, 0x7D, 0x44)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
BODY_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x8A, 0x8A, 0x8A)
VERY_LIGHT_BG = RGBColor(0xF5, 0xF3, 0xEF)
TABLE_HEADER_BG = RGBColor(0x2B, 0x5C, 0x8A)
TABLE_ROW_ALT = RGBColor(0xE8, 0xEE, 0xF4)
BORDER_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ═══════════════════════════════════════════════════════════════════
# FIRST: Rewrite all speaker notes in FIRST-PERSON perspective
# ═══════════════════════════════════════════════════════════════════

AUTHOR_NOTES = [
    # ── Slide 1: 封面 ──
    """【开场白 — 约30秒】

"老师好，我今天汇报的是我最近完成的一项研究工作。

题目是《环氧树脂微球包覆亲油性Fe₃O₄纳米颗粒作为示踪支撑剂用于分段产量监测》，简称ESP-T。

这项工作的核心思路是：用环氧树脂作为基体材料，把硬脂酸改性的纳米四氧化三铁包覆进去，做成一种既能支撑裂缝、又能长期监测各段产油量的双功能示踪支撑剂。

今天我想向您汇报四个方面的内容：研究背景与问题、实验方案与结果、核心发现与创新点、以及接下来的投稿计划。希望听取您的意见。" """,

    # ── Slide 2: 研究背景 ──
    """【研究背景 — 约1分30秒】

"首先介绍这项研究的出发点。

非常规油气（页岩油、致密油）目前占全球剩余可采储量的50%以上。水力压裂是开发这类资源的不可替代的核心技术。但压裂后有一个难题一直没解决得很好——不同压裂段的产油贡献差别很大，衰减速度也不一样。我们现在很难精确回答'每一段到底产了多少油'，导致后续的开发策略优化缺少数据支撑。这就是我做这项研究要解决的核心问题。

现有的示踪技术主要有三类，各有各的短板。传统油溶性化学示踪剂跟压裂液不兼容，没法长期监测。涂层型示踪支撑剂的致命缺陷是：涂层一溶解监测就终止了，而且制备要涂覆、固化、多次处理，工艺复杂。聚苯乙烯微球示踪剂虽然能监测油相，但PS的机械强度不够、热稳定性也差，200度以上就出问题了，深井根本用不了。

我梳理下来发现，行业的痛点很明确：缺一种能在井下长期稳定工作、专门针对油相的示踪方案。这正好为环氧树脂提供了一个切入点。" """,

    # ── Slide 3: 研究空白与创新点 ──
    """【切入点与创新 — 约1分钟】

"为什么选环氧树脂？这一点我想特别说明。

环氧树脂的机械强度、热稳定性和耐化学性都远超聚苯乙烯。之前Li 2021和Wei 2024两个团队已经在环氧树脂包覆水溶性示踪剂方面做了探索，但他们做的是水相——定位产水层段。我查了很多文献，目前确实没有人把环氧树脂作为亲油性释放基体来做油相产量监测。这是我的切入点。

具体来说，我的工作有三个层面的创新：

第一是材料设计层面——用乳液聚合一步实现低密度调控、纳米改性和原位示踪包覆，跟传统的多步骤涂覆路线完全不一样。

第二是表面工程层面——用硬脂酸做亲油改性，羧基跟Fe₃O₄表面羟基配位，长烷基链朝外构建疏水膜。这个设计的直接效果是把水接触角从72度拉到了104度，实现了亲水到疏水的翻转。

第三是模型方法层面——我建立了一个基于tanh平滑过渡的分段ADE模型，把示踪突破曲线分解成Gaussian脉冲分量和erfc拖尾分量，两个分量有明确的物理对应。最终R方做到了0.9939，这为之后的产量定量提供了数学基础。" """,

    # ── Slide 4: 实验方案 ──
    """【实验方案 — 约1分钟，快速过】

"实验方案我快速过一下。

合成路线分四步：第一步共沉淀法制备Mn掺杂的纳米Fe₃O₄；第二步硬脂酸乙醇溶液超声改性，得到亲油的nano-Fe₃O₄@SA；第三步是核心步骤——将改性纳米颗粒与E51环氧树脂、T31固化剂、中空玻璃微球预混，在SiO₂分散液和胍胶组成的稳定体系中进行乳液聚合；第四步50度固化、80度干燥。

表征我做了八个方向：SEM看形貌和元素分布，TGA/DSC看热稳定性，水接触角看润湿性变化，物理力学按行业标准SY/T 5107测试，油水过滤时间间接评价导流能力，ICP-MS定量示踪释放浓度，然后用K-P模型分析释放动力学机制，最后用分段ADE模型拟合突破曲线验证监测精度。整个链条从材料到性能到应用是打通的。" """,

    # ── Slide 5: SEM ──
    """【SEM结果 — 约1分钟】

"SEM是最直观的证据。

对比纯环氧微球和掺了nano-Fe₃O₄@SA的ESP-T，低倍下两种微球都呈现优异的球度和单分散性，说明乳液聚合的工艺可控性很好，加不加纳米填料不影响微球成型。

中高倍下差异就出来了。纯环氧表面是光滑的，ESP-T表面布满岛状凸起。这些凸起是硬脂酸包裹的纳米团簇，不是大块团聚——这说明硬脂酸改性确实有效抑制了相分离。

最关键的是高倍下一个细节：纳米团簇和环氧基体之间没有界面裂缝，呈现嵌入键合状态。这是硬脂酸长烷基链跟环氧分子链之间形成了物理缠结和疏水相互作用的直接证据。如果只是表面附着，界面一定会有裂缝或脱粘痕迹。

EDS面扫描也确认了——Fe元素信号遍布整个颗粒截面，不是只在表面。这证明纳米Fe₃O₄确实被包进了基体内部。

不过我需要在汇报中诚实地说：形成机制（固化收缩把纳米团簇挤出到表面）是我基于SEM结果的推断，还需要原位表征来确认动力学过程。这在论文里我也如实写了。" """,

    # ── Slide 6: 性能仪表盘 ──
    """【性能数据 — 约1分30秒，重点】

"材料做出来了，到底行不行？这一页是我们最硬的数据。

第一个指标：热稳定性。初始分解温度357.27度。典型井下温度是80到150度，深井最高也就200度。我的材料在200度以下只有微量水分挥发，环氧基体纹丝不动。这就是选环氧树脂、不选聚苯乙烯的根本原因——PS两百多度就扛不住了。

第二个指标：水接触角。从纯环氧的72.3度提升到104.6度，增加了32.3度。这个跃升意味着表面从弱亲水变成了明确疏水。原理是硬脂酸的羧基跟Fe₃O₄表面羟基配位成键，17个碳的长烷基链朝外定向排列，构建了一层疏水膜。

第三个指标：体积密度0.646克每立方厘米，比水轻。这意味着在水基压裂液中可以自然悬浮，不需要额外加悬浮剂。

下面的表格是按行业标准SY/T 5107测的：球度圆度超过0.9，酸溶解度3.3%远低于5%的标准限值，50MPa下破碎率2.9%，跟纯环氧的2.6%相当。这些硬指标说明ESP-T完全满足工业应用的基本门槛。" """,

    # ── Slide 7: 阻水亲油 ──
    """【阻水亲油 — 约1分钟】

"润湿性变化不只是接触角数字好看，它带来了实实在在的性能差异。

最直观的对比：油过滤时间从15分11秒降到了5分11秒，缩短了66%；而水过滤时间反而从不到3分钟延长到了28分钟41秒，增大了接近10倍。阻水亲油的效果非常明显。

这个设计的工程意义在于：在实际地层中，支撑剂充填层对油相提供增强导流，促进油向井筒流动；同时对水相增加阻力，帮助缓解水窜。这在含水率比较高的老井中可能会特别有价值。

需要注意的一点是，过滤时间是导流能力的间接评价，不是直接渗透率测量。如果要定量说导流能力提升了多少，需要做岩心驱替的渗透率测试。" """,

    # ── Slide 8: 释放动力学 ──
    """【释放动力学 — 约1分30秒】

"示踪剂怎么从支撑剂里释放出来？我用Korsmeyer-Peppas模型来定量描述。

K-P模型是一个幂律关系：C/C₀ = K·tⁿ。扩散指数n是最关键的参数——它直接告诉我们哪种机制在主导释放。球形载体中，n≤0.43是纯浓度梯度驱动的Fick扩散，n≥0.85是纯聚合物松弛驱动的Case-II传输，中间这个区间是两种机制协同。

我的数据：四个温度下n值全部落在0.45到0.85之间。这说明ESP-T的示踪释放不是简单的浓度差扩散，而是Fick扩散和聚合物松弛共同作用的结果。

速率常数K从30度的0.055提升到120度的0.196，系统性递增。这很好解释：温度越高，溶剂越容易渗透进环氧交联网络，溶胀越快，分子链间的约束力减弱，传输通道扩大。

四个温度的R方都在0.94以上，模型拟合质量没问题。但我需要主动说明一点：K-P模型是幂律形式，严格适用于释放初期(Mₜ/M∞ < 0.6)。我的14天数据C/C₀已经超过了1，所以模型不能外推到实验时限之外。这在论文里我也做了明确声明，审稿人应该不会在这点上找茬。" """,

    # ── Slide 9: ADE 模型 ⭐ ──
    """【核心结果 — 约2分钟，最重要的一页】

"这一页是我认为全文最重要的成果——用分段ADE模型拟合示踪突破曲线。

先交代一下为什么需要这个模型。示踪剂从支撑剂释放出来被油携带到井口，这个过程物理上天然分为两个阶段。关井期间积蓄的高浓度示踪剂团被一下子冲出来——这是ADE瞬时脉冲解，Gaussian形式的上升分量。残留在环氧基体里的示踪剂持续缓慢释放——这是ADE连续源解，erfc形式的拖尾分量。两个阶段之间我用tanh函数做平滑过渡，避免了阶跃拼接的数学不连续性。

现在看关键数据。

R方0.9939，RMSE 0.0210，全时间范围内残差都在正负两倍标准差以内。这个精度说明我的物理假设是正确的。

拟合流量0.46毫升每分钟，跟我泵设的0.5对比，相对误差只有8%。平均停留时间37.4分钟跟理论对流传输时间38.6分钟的比值是0.967。这两个验证说明模型参数不是纯经验拟合——它们有明确且正确的物理对应关系。

最后是最让我觉得有说服力的一个发现：我把C_rise和C_fall分别积分，发现erfc拖尾贡献了47%的总信号。将近一半！这意味着在流动条件下，基质扩散控制的持续释放仍然主导示踪传输。这从物理上证明了ESP-T做长期监测是可行的——不是所有示踪剂一股脑全出来了，而是有一半信号靠缓释在维持。

Peclet数约等于1也印证了这一点：对流和弥散旗鼓相当，恰好是一个过渡状态，跟基质扩散控制源的渐进释放特征完全自洽。" """,

    # ── Slide 10: 两相流 ──
    """【两相流验证 — 约1分钟】

"单相验证通过之后，我进一步测试了两相流条件下的表现。

实验设计了三组油水比和四个总流量水平。先看浓度数据：浓度随总流量增大而降低，这是稀释效应，单位体积流体接触时间短了。但浓度受油水比影响不大。

关键看示踪剂通量FO——单位时间通过井口的示踪剂质量，等价于支撑剂的释放速率。恒定油水比下FO跟总流量无关，因为油相跟支撑剂接触面积不变。但增大油水比，接触面积扩大，FO就上来了。

核心验证在右边：我把FO以单相油驱稳态值3.187微克每分钟为基准做归一化，然后跟实际油相流量对比——三个油水比下高度吻合。这证明了在稳态条件下，恒定总流量时可以从FO变化曲线反推标记段的实际产油量。

这是ESP-T能用于分段产量分配的实践基础。但这个结论有一个前提条件——稳态。如果生产中出现快速开井、关井、剧烈压力波动，FO跟实际流量的关系可能会偏离。这一点我在论文里也明确写了。" """,

    # ── Slide 11: 构效关系 ──
    """【构效关系 — 约45秒】

"这一页把整个工作的逻辑串起来。

从硬脂酸改性到长烷基链缠结，到纳米团簇嵌入键合，到疏水表面，到阻水亲油的导流特性，到非Fick释放机制，最后到ADE模型验证和分段产量监测。这是一条从微观分子设计到宏观工程应用的完整链条，每一环都有对应的实验证据支撑。

两个核心机制——润湿性转变和非Fick释放——不是孤立的，它们互相印证：疏水表面让油相更容易渗透进入环氧基体孔隙，而环氧基体本身的交联网络溶胀行为又决定了示踪剂的释放速率和模式。这两个机制的协同构成了ESP-T的科学基础。" """,

    # ── Slide 12: 局限性 ──
    """【局限性 — 约1分钟，主动提出】

"这部分我想诚实地汇报目前工作还存在哪些问题，也想听听您的意见。

第一，单段假设。目前实验和模型都是基于单裂缝段加均匀填充。实际矿场是多段、非均质的。从实验室到现场，这个可能是最大的鸿沟。

第二，稳态假设。FO标定依赖稳态条件。实际生产中的瞬变工况——开井、关井、快速降压——FO跟油相流量的关系可能需要重新标定。

第三，模型流体差异。我目前用的是十二烷模拟原油。真实原油里有沥青质、各种极性组分，粘度也不一样。沥青质会不会吸附到环氧表面影响释放？高粘度原油对溶胀动力学的影响有多大？这些我还没评估。

第四，工况边界。目前测试到120度。超过120度或者遇到高矿化度、含CO₂、H₂S的地层流体，环氧基体和硬脂酸改性层的长期稳定性还需要专门测试。

我想请老师特别指点一下：哪些局限是审稿人最可能在意的？我们应该在投稿前补充哪些实验或讨论？" """,

    # ── Slide 13: 总结 ──
    """【总结 — 约1分钟】

"最后总结一下我这篇工作的六点核心结论。

第一，成功制备了ESP-T。nano-Fe₃O₄@SA均匀包覆在环氧基体里，球度超过0.9，满足SY/T标准。

第二，性能上实现了三个关键指标：耐温357度远超井下需求、水接触角104.6度实现疏水、阻水亲油特性通过过滤时间实验得到了明确验证。

第三，释放机制清晰：用K-P模型确认了Fick扩散和Case-II松弛协同的非Fick异常传输。

第四，监测模型精度达到了我的预期：分段ADE模型R方0.9939，而且模型参数都有物理对应。

第五，两相流验证通过：示踪剂通量可以定量油相产量，为分段产量分配提供了实验依据。

第六，整体来看，这个材料把裂缝支撑和长期产量监测集成到了同一个体系里，在酸化压裂、深井、高压场景下应该都有应用潜力。

我想请老师帮我判断一下：这个工作整体的完整性够不够投一个好的期刊？还需要补充哪些实验或讨论？" """,
]

# ── Apply notes ────────────────────────────────────────────────────
for slide_idx, note_text in enumerate(AUTHOR_NOTES):
    if slide_idx < len(prs.slides):
        slide = prs.slides[slide_idx]
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.paragraphs[0].text = note_text
            for p in tf.paragraphs:
                p.font.size = Pt(10)
        except Exception as e:
            print(f"  WARNING slide {slide_idx+1}: {e}")

# ═══════════════════════════════════════════════════════════════════
# SECOND: Insert new slide after Slide 13 (conclusions) —
#         Journal Submission Plan
# ═══════════════════════════════════════════════════════════════════
# We need to insert a slide. python-pptx doesn't easily insert at position.
# Workaround: add a new blank slide, then move it via XML manipulation.

from lxml import etree

def add_new_slide_with_content(prs, slide_index):
    """Create a submission-plan slide and insert at position slide_index."""
    # Add slide at end first, then reorder
    layout = prs.slide_layouts[6]  # blank
    new_slide = prs.slides.add_slide(layout)

    # Build slide content
    # Title bar
    def add_rect(slide, left, top, width, height, fill_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(height))
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        return shape

    def add_textbox(slide, left, top, width, height, text="",
                    font_size=14, bold=False, color=BODY_TEXT,
                    alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.alignment = alignment
        return tf

    add_rect(new_slide, 0, 0, 13.333, 0.06, fill_color=ACCENT_BLUE)
    add_textbox(new_slide, 0.6, 0.25, 12, 0.55,
                "投稿计划与后续工作", font_size=26, bold=True, color=DARK_TEXT)
    add_textbox(new_slide, 0.6, 0.78, 12, 0.35,
                "拟投稿期刊比较  |  修稿计划  |  需要老师把关的方向",
                font_size=12, color=LIGHT_GRAY)
    # separator
    shape = new_slide.shapes.add_connector(1, Inches(0.6), Inches(1.15),
                                            Inches(12.7), Inches(1.15))
    shape.line.color.rgb = BORDER_LIGHT
    shape.line.width = Pt(0.5)

    # ── LEFT: Journal comparison table ──
    add_textbox(new_slide, 0.6, 1.4, 5, 0.35, "▎候选期刊对比", font_size=14, bold=True, color=ACCENT_BLUE)

    journal_data = [
        ["期刊", "IF (2024)", "分区", "审稿周期", "匹配度"],
        ["Fuel", "~7.5", "中科院1区/Top", "2-4月", "★★★★★"],
        ["Geoenergy Sci. Eng.", "~4.5", "中科院2区", "1-3月", "★★★★★"],
        ["Colloids Surf. A", "~5.2", "中科院2区", "1-2月", "★★★★"],
        ["ACS Omega", "~4.0", "中科院3区", "1-2月", "★★★★"],
        ["Energy & Fuels", "~5.5", "中科院2区/Top", "2-3月", "★★★"],
    ]
    table = new_slide.shapes.add_table(len(journal_data), 5,
                                        Inches(0.6), Inches(1.85),
                                        Inches(7.5), Inches(2.3)).table
    for r, row in enumerate(journal_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Microsoft YaHei"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                elif r == 1:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_ORANGE
                else:
                    p.font.color.rgb = BODY_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT

    # ── RIGHT TOP: Recommended strategy ──
    add_textbox(new_slide, 8.6, 1.4, 4.2, 0.35, "▎推荐策略", font_size=14, bold=True, color=ACCENT_ORANGE)
    add_rect(new_slide, 8.6, 1.85, 4.2, 2.3, fill_color=VERY_LIGHT_BG)

    strategy_text = (
        "首选: Fuel\n"
        "  Gong et al. 2024 同方向已发表\n"
        "  需强调: ADE模型 + 阻水亲油\n"
        "  为超越前人工作的主要亮点\n\n"
        "备选: Geoenergy Sci. Eng.\n"
        "  Zhao/Zhou/Li 的示踪支撑剂\n"
        "  均发表于此, 审稿人专业对口\n\n"
        "保底: ACS Omega / Colloids Surf. A\n"
        "  审稿快, 接受率较高"
    )
    add_textbox(new_slide, 8.8, 1.95, 3.8, 2.1, strategy_text, font_size=10, color=BODY_TEXT)

    # ── BOTTOM: Pre-submission checklist ──
    add_textbox(new_slide, 0.6, 4.45, 12, 0.35, "▎投稿前待办事项", font_size=14, bold=True, color=ACCENT_BLUE)

    checklist_items = [
        ("数据补充", "TGA数据是否需补充N₂气氛对比？\n120°C长期(>30天)稳定性数据？"),
        ("讨论深化", "与Gong 2024 PS微球的定量对比分析\n环氧vs PS综合性能雷达图"),
        ("语言润色", "建议投稿前请native speaker通读\n或使用专业润色服务"),
        ("格式适配", "目标期刊模板排版\n参考文献格式转换"),
        ("Cover Letter", "重点突出: 首次油相监测\nADE分段模型R²=0.9939"),
    ]
    for i, (title, desc) in enumerate(checklist_items):
        x = 0.6 + i * 2.5
        add_rect(new_slide, x, 4.9, 2.3, 2.0, fill_color=VERY_LIGHT_BG)
        add_rect(new_slide, x, 4.9, 2.3, 0.04, fill_color=ACCENT_GREEN)
        add_textbox(new_slide, x + 0.1, 5.0, 2.1, 0.3, title, font_size=11, bold=True, color=ACCENT_BLUE)
        add_textbox(new_slide, x + 0.1, 5.35, 2.1, 1.4, desc, font_size=9, color=BODY_TEXT)

    # ── Add speaker notes for this slide ──
    submission_note = """【投稿计划讨论 — 约2分钟】

"最后我想跟您讨论一下投稿计划。

我梳理了五个可能的目标期刊，列在这张表里。

首选我考虑投Fuel。IF大概7.5，中科院一区Top期刊。Gong et al. 2024那篇用PS微球做油相示踪支撑剂的文章就发在Fuel上，说明这个期刊认可这个方向。我们跟Gong的差异化很清楚：他用PS基体，我们用环氧树脂——耐温从200度拉到357度，机械强度也更好。再加上我们的ADE分段模型R方做到了0.9939、erfc拖尾47%的定量证据，这在方法学上比Gong更进一步。

备选是Geoenergy Science and Engineering，就是之前的Journal of Petroleum Science and Engineering。Zhao 2020、Zhou 2022、Li 2023这些做示踪支撑剂的文章全发在这个期刊上。好处是审稿人一定对口——他们审过这个细分方向的前人工作，能准确判断增量贡献。IF大概4.5，中科院二区。稳。

保底的话，ACS Omega审稿比较快一两个月，接受率相对高，但IF只有4左右而且是三区。Colloids and Surfaces A也差不多，Gong et al. 2024关于纳米Fe₃O₄示踪剂的另一篇文章就在那儿。

表格下面列了投稿前需要做的五件事，我想请老师帮我把关：

TGA数据现在只测了空气气氛，要不要补一个氮气惰性气氛做对比？还有120度长期超过30天的稳定性数据要不要补？

跟Gong 2024那篇文章的定量对比，我打算做一个综合性能雷达图，直观展示环氧vs PS在耐温、强度、密度、释放可控性这些维度上的差异。这个图我觉得放进论文会很有说服力。

语言方面，我写的是英文稿，投稿前还需要润色。

Cover Letter我想重点打两个亮点：首次用环氧树脂做油相示踪释放基体 + ADE分段模型的定量精度。

这些想听听您的意见，特别是首投的目标期刊选择和还需要补充的实验。" """

    try:
        ns = new_slide.notes_slide
        tf = ns.notes_text_frame
        tf.clear()
        tf.paragraphs[0].text = submission_note
        for p in tf.paragraphs:
            p.font.size = Pt(10)
    except:
        pass

    # ── Reorder: move last slide to position slide_index ──
    try:
        sldIdLst = prs.part.element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
        if sldIdLst is not None:
            sldId_elements = list(sldIdLst)
            last_el = sldId_elements[-1]
            sldIdLst.remove(last_el)
            if slide_index >= len(sldId_elements) - 1:
                sldIdLst.append(last_el)
            else:
                sldIdLst.insert(slide_index, last_el)
            print(f"  Inserted submission-plan slide at position {slide_index + 1}")
        else:
            print(f"  Added submission slide at end (reorder unavailable)")
    except Exception as e_reorder:
        print(f"  Added submission slide at end (reorder failed: {e_reorder})")

# Insert before the conclusions slide (which is at index 12)
add_new_slide_with_content(prs, 12)

# ── Save ───────────────────────────────────────────────────────────
prs.save(pptx_path)
print(f"\nSaved: {pptx_path}")
print(f"Total slides: {len(prs.slides)}")
print("All notes rewritten in first-person author perspective.")